from __future__ import annotations

import logging
from pathlib import Path
from typing import Sequence, Tuple

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from .errors import PPTBuildError
from .models import PPTBuildResult, SlideInfo

LOGGER = logging.getLogger(__name__)


class PPTBuilder:
    """Create PPT files from extracted slide images."""

    def __init__(
        self,
        fill_mode: bool = True,
        margin_inches: float = 0.3,
    ) -> None:
        self.fill_mode = fill_mode
        self.margin_inches = max(float(margin_inches), 0.0)
        # 16:9 比例 (标准1920x1080): 10 x 5.625 inches
        self.slide_width = Inches(10)
        self.slide_height = Inches(5.625)

    def build(
        self,
        slides: Sequence[SlideInfo],
        output_path: Path,
        title: str | None = None,
        subtitle: str | None = None,
    ) -> PPTBuildResult:
        if not slides:
            raise PPTBuildError("No slides to add into PPT.")

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        presentation = Presentation()
        presentation.slide_width = self.slide_width
        presentation.slide_height = self.slide_height

        if title:
            self._add_title_slide(presentation, title, subtitle)

        for slide_info in slides:
            self._add_image_slide(presentation, slide_info)

        presentation.save(str(output_path))
        LOGGER.info("PPT generated: %s (slides=%s)", output_path, len(slides))
        return PPTBuildResult(ppt_path=output_path, slide_count=len(slides))

    def _add_title_slide(self, presentation: Presentation, title: str, subtitle: str | None) -> None:
        layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        if subtitle:
            try:
                subtitle_shape = slide.placeholders[1]
            except IndexError:
                subtitle_shape = None
            if subtitle_shape:
                subtitle_shape.text = subtitle

    def _add_image_slide(self, presentation: Presentation, slide_info: SlideInfo) -> None:
        layout = presentation.slide_layouts[6]  # Blank layout
        slide = presentation.slides.add_slide(layout)
        left, top, width, height = self._calculate_bounds(slide_info.path)
        slide.shapes.add_picture(str(slide_info.path), left, top, width=width, height=height)

    def _calculate_bounds(self, image_path: Path) -> Tuple[int, int, int, int]:
        with Image.open(image_path) as img:
            img_width, img_height = img.size

        margin = Inches(self.margin_inches)
        available_width = max(self.slide_width - 2 * margin, Inches(0.01))
        available_height = max(self.slide_height - 2 * margin, Inches(0.01))

        width_ratio = available_width / img_width
        height_ratio = available_height / img_height

        if self.fill_mode:
            scale_ratio = max(width_ratio, height_ratio)
        else:
            scale_ratio = min(width_ratio, height_ratio)

        display_width = int(img_width * scale_ratio)
        display_height = int(img_height * scale_ratio)

        left = int((self.slide_width - display_width) / 2)
        top = int((self.slide_height - display_height) / 2)

        return left, top, display_width, display_height
